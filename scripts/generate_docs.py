import markdown
import pdfkit
from pptx import Presentation
import os

def md_to_pdf(md_file):
    with open(md_file, 'r', encoding='utf-8') as f:
        html = markdown.markdown(f.read(), extensions=['tables'])
    pdf_file = md_file.replace('.md', '.pdf')
    pdfkit.from_string(html, pdf_file)
    print(f'PDF 생성 완료: {pdf_file}')

def md_to_pptx(md_file, pptx_file):
    prs = Presentation()
    with open(md_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    slide = None
    for line in lines:
        if line.startswith('#'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = line.strip('#').strip()
        elif slide:
            textbox = slide.shapes.placeholders[1]
            textbox.text += line
    prs.save(pptx_file)
    print(f'슬라이드 생성 완료: {pptx_file}')

if __name__ == '__main__':
    md_files = [
        'docs/ADMIN_GUIDE.md',
        'docs/OPERATION_SECURITY_CHECKLIST.md',
        'docs/OPERATION_WORKFLOW.md'
    ]
    for md_file in md_files:
        if os.path.exists(md_file):
            md_to_pdf(md_file)
            pptx_file = md_file.replace('.md', '.pptx')
            md_to_pptx(md_file, pptx_file) 